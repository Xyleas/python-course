# 1. Import Library
from pptx import Presentation

# 2. Create presentation
presentation = Presentation()

# 3. Create slide
# 3.1. Layout selection
slide_layout = presentation.slide_layouts[1] # Title and content layout
# 3.2 . Slide creation
slide = presentation.slides.add_slide(slide_layout)

# 4. Add title
slide.placeholders[0].text = 'My Title'

# 5. Add content
slide.placeholders[1] = 'Hello World\nBye'

# 6. Save presentation
presentation.save("test.pptx")