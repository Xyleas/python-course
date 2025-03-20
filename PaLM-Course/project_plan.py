import google as genai
from openpyxl import Workbook
from pptx import Presentation

def generate(api_key, project_charter):
    genai.configure(api_key=api_key)

    # Create the model
    generation_config = {
    "temperature": 1,
    "top_p": 0.95,
    "top_k": 40,
    "max_output_tokens": 8192,
    "response_mime_type": "text/plain",
    }

    model = genai.GenerativeModel(
    model_name="gemini-2.0-flash-exp",
    generation_config=generation_config,
    )

    chat_session = model.start_chat(
    history=[
        {
        "role": "user",
        "parts": [
            "Act as a senior project manager with experience in software products. Write a Project plan based on the Project Charter below:\n\n{{input}}\n\n--\nEND OF PROJECT CHARTER\n\nThe output of the Project plan needs to be in tabular format. The columns are:\n1) Take Name: Description of the tasks.\n2) Duration: Time required to complete the task.\n3) Dependencies: Other tasks that must be completed before this task.\n4) Status: Current status of the task (e.g., Not Started, In Progress, Completed).\n5) Resources: Tools, team members, software, infrastructure, etc., required for the task.\n\nExample of the output:\n\nTask Name | Duration | Dependencies | Status | Resources\nCreating Plan | 1 Day | None | Not Started | Leadership team\nDocumentation | 3 days | Creating Plan | Not Started | Development team, internal software\n\nDon't add any additional content or notes after you finish listing the tasks. Add at least 10 tasks and no extra content after. Now please write the Project Plan in tabular format:\n",
        ],
        },
        {
        "role": "model",
        "parts": [
            "| Task Name | Duration | Dependencies | Status | Resources |\n|---|---|---|---|---|\n| Project Kick-off Meeting | 1 Day | None | Not Started | Project Manager, Stakeholders, Core Team |\n| Detailed Requirements Gathering | 5 Days | Project Kick-off Meeting | Not Started | Business Analyst, Product Owner, Stakeholders |\n| System Architecture Design | 7 Days | Detailed Requirements Gathering | Not Started | Lead Developer, Architects |\n| UI/UX Design | 7 Days | Detailed Requirements Gathering | Not Started | UI/UX Designers |\n| Database Design | 3 Days | System Architecture Design | Not Started | Database Administrator, Lead Developer |\n| Development - Core Functionality | 20 Days | System Architecture Design, Database Design | Not Started | Development Team |\n| Development - UI Implementation | 15 Days | UI/UX Design, Development - Core Functionality | Not Started | Front-end Developers |\n| Testing - Functional Testing | 10 Days | Development - UI Implementation | Not Started | QA Testers, Development Team |\n| Testing - Performance & Security Testing | 5 Days | Development - UI Implementation | Not Started | QA Testers, Security Specialists |\n| User Acceptance Testing (UAT) | 5 Days | Testing - Functional Testing, Testing - Performance & Security Testing | Not Started |  Product Owner, Stakeholders, End Users |\n| Deployment & Release | 2 Days | User Acceptance Testing (UAT) | Not Started | DevOps Team, Project Manager |\n| Post-Launch Monitoring | Ongoing | Deployment & Release | Not Started | DevOps Team, Support Team |\n",
        ],
        },
    ]
    )

    response = chat_session.send_message("INSERT_INPUT_HERE")

    print(response.text)

    rows = tabular_data.split('\n')

    spreadsheet = []

    for row in rows:
        split_row = row.split('|')
        spreadsheet.append(split_row)
    
    #testing print(rows)
    print(spreadsheet)

    excel_save(spreadsheet)
    powerpoint_save(spreadsheet)

    return response.result

def excel_save(spreadsheet):
    wb = Workbook()
    ws = wb.active

    for row in spreadsheet:
        ws.append(row)

    wb.save('project_plan.xlsx')

def powerpoint_save(spreadsheet):
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[1] # Title and content layout

    spreadsheet = spreadsheet[2:]

    for row in spreadsheet:
        if len(row) < 5:
            continue

        slide = presentation.slides.add_slide(slide_layout)

        # Title
        slide.placeholders[0].text = row[0]

        # Content
        slide.placeholders[1].text = "Duration: " + row[1] + "\nDependencies: " + row[2] + "\nStatus: " + row[3] + "\nResources: " + row[4]

    presentation.save('project_plan.pptx')